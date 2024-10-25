import os
from glob import glob

import extract_msg
import pandas as pd
import pptx
import PyPDF2
import textract
from docx import Document

## This notebook contains all the relevant functions used for the word counting for woo request mix of reading and counting.
## Author: Michael de Winter


# Different functions for opening files.
def read_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        pdf_reader = PyPDF2.PdfReader(file)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
    return text


def read_docx(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)


def read_msg(file_path):
    msg = extract_msg.Message(file_path)
    msg_text = "Sender: " + str(msg.sender) + " | \n "
    msg_text = msg_text + "To: " + str(msg.to) + " | \n "
    msg_text = msg_text + "CC: " + str(msg.cc) + " | \n "
    msg_text = msg_text + "BCC: " + str(msg.bcc) + " | \n "
    msg_text = msg_text + "Subject: " + str(msg.subject) + " | \n "
    msg_text = msg_text + "Body: " + str(msg.body)
    return msg_text


def read_xls(file_path):
    # Convert the pandas DataFrame to a single string
    data = pd.read_excel(file_path)
    collapsed_text = data.to_string(index=False)
    return collapsed_text


def read_doc(file_path):
    text = textract.process(file_path)
    return text


def read_pptx(file_path):
    text = ""
    presentation = pptx.Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text


def count_word_combinations(text, word_list):
    """
    Count the number of occurences of words in a string text.

    TODO: This could be included in further tests
    @param text: string value of text.
    @param word_list: a array of words which have to found in the text.
    """
    counts = {}

    # Convert the text to lower case for case-insensitive matching
    text = str(text.lower())

    # Iterate through each combination of words
    for i in range(len(word_list)):
        # Count the number of occurrences of the combination in the text
        count = text.count(word_list[i])
        # Store the count in the dictionary
        counts[word_list[i]] = count
    return counts


def process_file(a_file_path):
    """
    Open a file  based on a file extension.

    @param a_file_path: File path to a file.
    """

    file_split = a_file_path.split(".")
    file_type = file_split[-1].lower()
    a_content = ""

    if len(file_split) > 1:
        try:
            if "pdf" in file_type:
                # Action for PDF file type
                a_content = read_pdf(a_file_path)
            elif "xls" in file_type or "xlsx" in file_type:
                # Action for XLS file type
                a_content = read_xls(a_file_path)
            elif "docx" in file_type:
                # Action for DOCX file type
                a_content = read_docx(a_file_path)
            elif "doc" in file_type:
                a_content = read_doc(a_file_path)
            elif "msg" in file_type:
                # Action for MSG file type
                a_content = read_msg(a_file_path)
            elif "txt" in file_type:
                with open(a_file_path, "r") as file:
                    # Read the content of the file
                    a_content = file.read()

            elif "pptx" in file_type:
                a_content = read_pptx(a_file_path)
            else:
                a_content = textract.process(a_file_path)
        except Exception as e:
            print(e)

    return a_content


def word_count_in_file(a_file_path, aword_list, verbose_setting=False):
    """
    Counts the number of words in a text based on a word_list.

    TODO: This can be used in tests for further validation.
    @param a_file_path: File path to a file.
    @param aword_list: A list of words which should be counted
    """

    return_words = count_word_combinations(process_file(a_file_path), aword_list)

    return return_words


def count_files_in_directory(directory):
    file_count = 0
    for root, dirs, files in os.walk(directory):
        file_count += len(files)
    return file_count


def retrieve_counts_extension_types(adirectory):
    """Looks in a directory to find the number of files of each type.

    TODO: This function can be included in tests
    @param adirectory: a statement of a directory which can be read by glob.
    """
    file_types = {}

    for file in glob(adirectory):
        file_type = file.split(".")[-1]
        if file_type not in file_types:
            file_types[file_type] = 1
        else:
            file_types[file_type] += 1

    # Sort the resulting list by file type in ascending order
    sorted_file_types = sorted(file_types.items(), key=lambda x: x[1])

    # Print the sorted file types
    return sorted_file_types


class FilterMinWords:
    """
    This class is used to filter out certain documents if that contain words in title or the filepath.

    """

    def __init__(self, afilename, min_words_path, aproject):
        if isinstance(afilename, pd.DataFrame):
            self.hits = afilename
        # Read in the file if it's only a file path
        elif ".xls" in afilename or ".xlsx" in afilename:
            self.hits = pd.read_excel(afilename)

        self.mins = pd.read_excel(min_words_path, engine="odf", header=None)
        self.project = aproject

    def filter(self):
        "Filter out in filepath and in filename"
        if self.project == "woo_shell_2024":
            for word in self.mins[0][0:56]:
                self.hits = self.hits[
                    ~self.hits["Locatie in iDMS"].str.contains(word.lower())
                ]
                # self.hits =self.hits[~self.hits['Bestandsnaam in export'].str.lower().str.contains(word.lower())]

        return self.hits
