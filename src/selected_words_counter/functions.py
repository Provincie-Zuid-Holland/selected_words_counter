import os
import shutil
from datetime import datetime
from glob import glob

import extract_msg
import pandas as pd
import pptx
import PyPDF2
import textract
from docx import Document

import config

""" 
This code contains all functions for reading in different file formats.


@author: Michael de Winter
"""


# Different functions for opening files.
def read_pdf(file_path):
    text = ""
    try:
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()

                # Handle encoding issues with utf-8-sig and replace errors
                if page_text:
                    text += page_text.encode("utf-8-sig", errors="replace").decode(
                        "utf-8-sig", errors="replace"
                    )

    except Exception as e:
        print(f"Error processing {file_path}: {e}")

    return text


def read_docx(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return "\n".join(text)


def read_msg(file_path):
    msg = extract_msg.Message(file_path)
    msg_text = "Sender: " + str(msg.sender) + " | \n "
    msg_text = msg_text + "To: " + str(msg.to) + " | \n "
    msg_text = msg_text + "CC: " + str(msg.cc) + " | \n "
    msg_text = msg_text + "BCC: " + str(msg.bcc) + " | \n "
    msg_text = msg_text + "Subject: " + str(msg.subject) + " | \n "
    msg_text = msg_text + "Body: " + str(msg.body)
    return msg_text


def read_xls(file_path):
    # Convert the pandas DataFrame to a single string
    data = pd.read_excel(file_path)
    collapsed_text = data.to_string(index=False)
    return collapsed_text


def read_doc(file_path):
    text = textract.process(file_path)
    return text


def read_pptx(file_path):
    text = ""
    presentation = pptx.Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text


def process_file(a_file_path):
    """

    Open a file based on a file extension.

    @param a_file_path: File path to a file.
    """

    file_split = a_file_path.split(".")
    file_type = a_file_path.rsplit(".", 1)[-1]
    a_content = ""

    if len(file_split) > 1:
        try:
            if "pdf" in file_type:
                # Action for PDF file type
                a_content = read_pdf(a_file_path)
            elif "xls" in file_type or "xlsx" in file_type:
                # Action for XLS file type
                a_content = read_xls(a_file_path)
            elif "docx" in file_type:
                # Action for DOCX file type
                a_content = read_docx(a_file_path)
            elif "doc" in file_type:
                a_content = read_doc(a_file_path)
            elif "msg" in file_type:
                # Action for MSG file type
                a_content = read_msg(a_file_path)
            elif "txt" in file_type:
                with open(a_file_path, "r") as file:
                    # Read the content of the file
                    a_content = file.read()
            elif "pptx" in file_type:
                a_content = read_pptx(a_file_path)
        except Exception as e:
            print("Errow with " + a_file_path)
            print(e)

    return a_content


def generate_filename(output_dir, version):
    """
    Generate a filename output based on the current date and the version.

    @param output_dir: The directory to base the filename on.
    @param version: Give a version to the output name.
    """
    current_date = datetime.now().strftime("%Y-%m-%d")  # Format as YYYY-MM-DD
    filename = f"{output_dir}/selected_word_counter_output_{current_date}_v{version}"
    return filename


class FilterMinWords:
    """
    This class is used to filter out certain documents if that contain words in title or the filepath.

    """

    def __init__(self, afilename, min_words_path, aproject):
        if isinstance(afilename, pd.DataFrame):
            self.hits = afilename
        # Read in the file if it's only a file path
        elif ".xls" in afilename or ".xlsx" in afilename:
            self.hits = pd.read_excel(afilename)

        self.mins = pd.read_excel(min_words_path, engine="odf", header=None)
        self.project = aproject

    def filter(self):
        "Filter out in filepath and in filename"
        if self.project == "woo_shell_2024":
            for word in self.mins[0][0:56]:
                self.hits = self.hits[
                    ~self.hits["Locatie in iDMS"].str.contains(word.lower())
                ]
                # self.hits =self.hits[~self.hits['Bestandsnaam in export'].str.lower().str.contains(word.lower())]

        return self.hits


def replace_last_slash(path, replacement=""):
    """
    This function is used for testing.

    @param path: path to be split.
    @param replacement: value to replace / with.
    """
    parts = path.rsplit("/", 1)
    return replacement.join(parts)


def delete_directory(path):
    try:
        shutil.rmtree(path)
        print(f"Directory '{path}' and all its contents were deleted successfully.")
    except FileNotFoundError:
        print(f"Directory '{path}' not found.")
    except Exception as e:
        print(f"An error occurred: {e}")
