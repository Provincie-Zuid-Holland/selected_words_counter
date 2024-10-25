import os
import random
import string
import zipfile

import openpyxl
import requests
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas

try:
    import win32com.client
except ImportError:
    win32com = (
        None  # win32com is only available on Windows systems with pywin32 installed
    )

word_site = "https://www.mit.edu/~ecprice/wordlist.10000"

response = requests.get(word_site)
WORDS = [line.decode("utf-8") for line in response.content.splitlines()]


def random_text(length):
    """Generate a random string of fixed length."""
    return " ".join(random.choice(WORDS) for _ in range(length))


def create_docx(filename):
    """Create a synthetic .docx file."""
    doc = Document()
    for i in range(random.randint(0, 50)):
        doc.add_heading(f"{random.choice(WORDS)} {i+1}", level=1)
        doc.add_paragraph(random_text(200))
    doc.save(filename)
    print(f"Created {filename}")


def create_pptx(filename):
    """Create a synthetic .pptx file."""
    prs = Presentation()
    for i in range(random.randint(0, 50)):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"{random.choice(WORDS)} {i+1}"
        subtitle.text = random_text(100)
    prs.save(filename)
    print(f"Created {filename}")


def create_xlsx(filename):
    """Create a synthetic .xlsx file."""
    wb = openpyxl.Workbook()
    ws = wb.active
    for row in range(1, 11):
        for col in range(1, 6):
            ws.cell(row=row, column=col).value = random.randint(1, 100)
    wb.save(filename)
    print(f"Created {filename}")


def create_pdf(filename):
    """Create a synthetic .pdf file."""
    c = canvas.Canvas(filename)
    for i in range(5):
        text = random_text(100)
        c.drawString(100, 800 - i * 150, text)
    c.save()
    print(f"Created {filename}")


def create_zip(filename):
    """Create a synthetic .zip file containing random text files."""
    with zipfile.ZipFile(filename, "w") as zipf:
        for i in range(3):
            temp_filename = f"temp_file_{i}.txt"
            with open(temp_filename, "w") as f:
                f.write(random_text(500))
            zipf.write(temp_filename)
            os.remove(temp_filename)
    print(f"Created {filename}")


def create_msg(filename):
    """Create a synthetic .msg file (requires Windows and Outlook)."""
    if win32com is None:
        print("win32com module not available. Cannot create .msg files.")
        return
    try:
        outlook = win32com.client.Dispatch("Outlook.Application")
        mail = outlook.CreateItem(0)
        mail.To = "recipient@example.com"
        mail.Subject = "Synthetic Email"
        mail.Body = random_text(500)
        # Save the email as a .msg file
        mail.SaveAs(os.path.abspath(filename))
        print(f"Created {filename}")
    except Exception as e:
        print(f"Could not create {filename}: {e}")
        print(
            "Ensure that you are running this on Windows with Outlook installed and pywin32 package."
        )


def main():
    # Params for the amount of documents made.
    amount_of_documents = 70
    for xi in range(0, amount_of_documents):
        if random.randint(0, 50) > 10:
            create_docx("./data/" + random_text(6) + ".docx")
        if random.randint(0, 50) > 25:
            create_pptx("./data/" + random_text(6) + ".pptx")
        if random.randint(0, 50) > 35:
            create_xlsx("./data/" + random_text(6) + ".xlsx")
        if random.randint(0, 50) > 15:
            create_pdf("./data/" + random_text(6) + ".pdf")
        if random.randint(0, 50) > 45:
            create_zip("./data/" + random_text(6) + ".zip")
        if random.randint(0, 50) > 45:
            create_msg("./data/" + random_text(6) + ".msg")


if __name__ == "__main__":
    main()
